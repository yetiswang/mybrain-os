#!/usr/bin/env python3
"""
Form entry -> PowerPoint profile card.

Turns a structured entry (e.g. an onboarding-form submission: name, role,
contact, photo, preferences) into a new slide in a directory deck, by
*mirroring an existing template slide* rather than building a layout from
scratch. The deck stays visually consistent because every card is a copy of
the same hand-designed template with only the text, toggles, and photo swapped.

Why mirror instead of build:
  - python-pptx can't clone a whole slide cleanly, but it CAN deep-copy the
    XML of individual shapes. Copy every shape from a template slide into a
    fresh slide, then edit text runs in place -> the design (fonts, colours,
    positions) is preserved exactly, for free.
  - Skip the template's picture; add the entry's photo fresh, clipped to a
    circle, at the template photo's locked position so all cards align.

The clipped-circle photo is the fiddly part — see crop_square() and the
prstGeom 'ellipse' trick below.

Genericised from a real "team directory" deck generator. Replace the field
map and shape ids with your own template's.
"""
import copy
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn

DECK = os.path.expanduser("<directory-deck>.pptx")
TEMPLATE_SLIDE = 5      # 0-based index of the slide to mirror (an existing card)
INSERT_BEFORE = 8       # insert new card before this slide index (e.g. before "add yours" templates)

# Photo position, locked so every card aligns (inches): left, top, side
PHOTO_L, PHOTO_T, PHOTO_S = 0.28, 1.20, 1.95

# Map your template's shape ids -> which entry field fills them.
# Find ids with: for sh in slide.shapes: print(sh.shape_id, sh.name, sh.text_frame.text)
FIELD_SHAPE_IDS = {
    "name": 8,
    "role": 9,
    "quote": 30,
    "field_a": 32,        # e.g. "area of work"
    "field_b": 34,        # e.g. "expertise"
    "find_me": 36,
}

# Contact-preference pills: shape name -> the preference key it represents.
PILL_SHAPES = {
    "Rounded Rectangle 37": "email",
    "Rounded Rectangle 38": "phone",
    "Rounded Rectangle 39": "online",
    "Rounded Rectangle 40": "in_person",
}
TEAL = RGBColor(0x00, 0x80, 0x80)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SOFT = RGBColor(0x78, 0x8C, 0x91)
HAIRLINE = RGBColor(0xCF, 0xDC, 0xDC)


def crop_square(src, out, face_frac_x=0.40):
    """Crop a face-centred square. face_frac_x = horizontal centre of the face
    as a fraction of width (portraits are rarely centred)."""
    from PIL import Image
    im = Image.open(src).convert("RGB")
    w, h = im.size
    side = min(w, h)
    left = int(min(max(0, face_frac_x * w - side / 2), w - side))
    im.crop((left, 0, left + side, side)).save(out, quality=92)
    return out


def set_text(shape, text):
    """Replace the first run's text, preserving its formatting."""
    shape.text_frame.paragraphs[0].runs[0].text = text


def set_pill(shape, on):
    r = shape.text_frame.paragraphs[0].runs[0]
    if on:
        shape.fill.solid(); shape.fill.fore_color.rgb = TEAL
        shape.line.fill.background()
        r.font.bold = True; r.font.color.rgb = WHITE
    else:
        shape.fill.solid(); shape.fill.fore_color.rgb = WHITE
        shape.line.color.rgb = HAIRLINE; shape.line.width = Emu(12700)
        r.font.bold = False; r.font.color.rgb = SOFT


def add_clipped_circle(slide, photo):
    pic = slide.shapes.add_picture(photo, Inches(PHOTO_L), Inches(PHOTO_T),
                                   Inches(PHOTO_S), Inches(PHOTO_S))
    spPr = pic._element.spPr
    for tag in ("a:prstGeom", "a:custGeom"):       # drop any existing geometry
        el = spPr.find(qn(tag))
        if el is not None:
            spPr.remove(el)
    # keep <a:xfrm> first, then append the ellipse preset -> picture clips to a circle
    prst = spPr.makeelement(qn("a:prstGeom"), {"prst": "ellipse"})
    prst.append(prst.makeelement(qn("a:avLst"), {}))
    spPr.append(prst)
    pic.line.fill.background()                     # no ring; clean edge


def add_card(entry):
    prs = Presentation(DECK)
    src = prs.slides[TEMPLATE_SLIDE]
    new = prs.slides.add_slide(src.slide_layout)

    # blank the new slide, then deep-copy every shape from the template EXCEPT the picture
    for sh in list(new.shapes):
        sh._element.getparent().remove(sh._element)
    for sh in src.shapes:
        if sh.shape_type == 13:        # PICTURE -> add a fresh one later
            continue
        new.shapes._spTree.append(copy.deepcopy(sh._element))

    by_id = {sh.shape_id: sh for sh in new.shapes}
    for field, sid in FIELD_SHAPE_IDS.items():
        if entry.get(field):
            set_text(by_id[sid], entry[field])

    for sh in new.shapes:
        if sh.name in PILL_SHAPES:
            set_pill(sh, PILL_SHAPES[sh.name] in entry.get("contact_prefs", []))

    if entry.get("photo"):
        square = crop_square(entry["photo"], "/tmp/_card_square.jpg")
        add_clipped_circle(new, square)

    # move the new (currently last) slide to the desired position
    lst = prs.slides._sldIdLst
    sid = list(lst)[-1]
    lst.remove(sid); lst.insert(INSERT_BEFORE, sid)

    prs.save(DECK)
    print("added card for %r at slide index %d" % (entry.get("name"), INSERT_BEFORE))


if __name__ == "__main__":
    # Example entry — wire this to your form / email parser.
    add_card({
        "name": "Jordan Lee",
        "role": "Lab Manager · Some Department",
        "quote": "“Better safe than sorry”",
        "field_a": "Wet lab · imaging · sample prep",
        "field_b": "Lab safety, procurement, onboarding.",
        "find_me": "Building X, room 1.23",
        "contact_prefs": ["email", "online", "in_person"],
        "photo": "<photo>.jpg",
    })
