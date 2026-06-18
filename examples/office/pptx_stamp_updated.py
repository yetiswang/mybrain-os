#!/usr/bin/env python3
"""
Stamp an "Updated <date>, <time>" tag on a deck's title slide.

Call this as the first step of your deck-save workflow (before regenerating
a PDF and copying backups), so the title slide / PDF page 1 always shows when
the deck last changed. Idempotent: it removes any prior "Updated …" tag before
adding the new one, so re-running never stacks duplicates.

Pairs naturally with pptx_card_from_form.py — bump the stamp on every edit.
"""
import os
import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

DECK = os.path.expanduser("<directory-deck>.pptx")
TITLE_SLIDE = 0
TAG_PREFIX = "Updated "
SOFT = RGBColor(0x78, 0x8C, 0x91)


def stamp_text():
    now = datetime.datetime.now().astimezone()
    # e.g. "Updated 18 June 2026, 16:22 CEST"
    return "%s%s %s, %s %s" % (TAG_PREFIX, now.day, now.strftime("%B %Y"),
                               now.strftime("%H:%M"), now.strftime("%Z"))


def main():
    prs = Presentation(DECK)
    s = prs.slides[TITLE_SLIDE]
    for sh in list(s.shapes):                       # remove prior tag (idempotent)
        if sh.has_text_frame and sh.text_frame.text.strip().startswith(TAG_PREFIX):
            sh._element.getparent().remove(sh._element)
    # top-right, balancing a logo on the left; tune to your title layout
    tb = s.shapes.add_textbox(Inches(5.15), Inches(0.66), Inches(4.5), Inches(0.3))
    tf = tb.text_frame
    tf.word_wrap = False
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    r = p.add_run(); r.text = stamp_text()
    r.font.name = "Calibri"; r.font.size = Pt(11); r.font.color.rgb = SOFT
    prs.save(DECK)
    print(stamp_text())


if __name__ == "__main__":
    main()
